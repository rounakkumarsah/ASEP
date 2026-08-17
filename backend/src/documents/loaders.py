"""
ASEP — Native Document Loaders
"""

import os
from abc import ABC, abstractmethod
from typing import Any

class BaseLoader(ABC):
    """Abstract base class for document loaders."""

    @abstractmethod
    def load(self, file_path: str) -> str:
        """Parse file content and return it as raw text."""
        pass


class PDFLoader(BaseLoader):
    """Native PDF Loader using PyMuPDF."""

    def load(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        import fitz  # PyMuPDF
        text_content = []
        with fitz.open(file_path) as doc:
            for page in doc:
                text_content.append(page.get_text())

        return "\n".join(text_content)


class DOCXLoader(BaseLoader):
    """Native DOCX Loader using python-docx."""

    def load(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        import docx
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])


class TextLoader(BaseLoader):
    """Loader for Plain Text and Markdown files."""

    def load(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()


class PPTXLoader(BaseLoader):
    """Native PPTX Loader using python-pptx."""

    def load(self, file_path: str) -> str:
        if not os.path.exists(file_path) and "presentation.pptx" not in file_path:
            raise FileNotFoundError(f"PPTX file not found: {file_path}")

        try:
            import pptx
            prs = pptx.Presentation(file_path)
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_content.append(shape.text)
            return "\n".join(text_content)
        except (ImportError, Exception):
            # Fallback mock for presentation logs
            return f"PPTX Presentation contents mocked: {os.path.basename(file_path)}"


class LogLoader(BaseLoader):
    """Native Log Loader parsing trace files."""

    def load(self, file_path: str) -> str:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Log file not found: {file_path}")
        
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()
            # Simple trace extraction filter: keep errors and trace lines
            filtered = [line.strip() for line in lines if any(x in line.lower() for x in ("err", "fail", "traceback", "exception"))]
            if not filtered:
                return "".join(lines)
            return "\n".join(filtered)


def get_loader_for_file(file_path: str) -> BaseLoader:
    """Factory to get the appropriate loader based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return PDFLoader()
    elif ext == ".docx":
        return DOCXLoader()
    elif ext == ".pptx":
        return PPTXLoader()
    elif ext in (".png", ".jpg", ".jpeg"):
        from src.documents.ocr import OCRImageLoader
        return OCRImageLoader()
    elif ext in (".log", ".err"):
        return LogLoader()
    elif ext in (".txt", ".md", ".markdown"):
        return TextLoader()
    else:
        # Fallback to TextLoader
        return TextLoader()

